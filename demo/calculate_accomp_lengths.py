from pptx import Presentation
from pptx.util import Inches, Pt
import csv
import os
import sys, getopt
from subprocess import Popen, PIPE, STDOUT
from time import sleep
import pyautogui
from pyautogui import hotkey

def usage():
    print('Usage: calculate_accomp_lengths.py -i <inputfile> -f <fontname> -s <fontsize>')

def main():
    try:
        opts, args = getopt.getopt(sys.argv[1:],"hi:r:f:s:",["help","input=","range=","font=","fontsize="])
    except getopt.GetoptError as err:
        print('Unrecognized input')
        print(err)
        usage()
        sys.exit(2)
    csv_filename = None
    ppt_output_filename = None
    font_name = None
    font_size = None
    for o, a in opts:
        if o in ("-h", "--help"):
            usage()
            sys.exit()
        elif o in ("-i", "--input"):
            csv_filename = a
            print("Input: " + csv_filename)
            output_filename = os.path.splitext(csv_filename)[0]+"_accomp_lengths.csv"
            print("Output: " + output_filename)
            ppt_output_filename = os.path.splitext(csv_filename)[0]+"_accomp_test.pptx"
        elif o in ("-r", "--irange"):
            a = a.split(',')
            accomp_min = int(a[0])
            accomp_max = int(a[1])
            print("Using range: [" + str(accomp_min) + "," + str(accomp_max) + "]")
        elif o in ("-f", "--font"):
            font_name = a
            print("Using font: " + font_name)
        elif o in ("-s", "--fontsize"):
            font_size = int(a)
            print("Using font size: " + str(font_size))
        else:
            assert False, "unhandled option"

    if csv_filename == None:
        usage()
        sys.exit()

    os.system('color 0')
    class bcolors:
        HEADER = '\033[95m'
        OKBLUE = '\033[96m'
        OKGREEN = '\033[92m'
        WARNING = '\033[93m'
        FAIL = '\033[91m'
        ENDC = '\033[0m'
        BOLD = '\033[1m'
        UNDERLINE = '\033[4m'

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    word_list = []

    with open(csv_filename, 'r', encoding='UTF-8', newline='') as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=',')
        line_count = 0
        left = Inches(0)
        top = Inches(0)
        # this size doesn't matter, just long enough to not wrap
        width = Inches(20)
        height = Inches(5)

        for row in csv_reader:
            if line_count == 0:
                print(bcolors.WARNING + 'Imported CSV Columns:' + bcolors.ENDC)
                print(', '.join(row)) 
                line_count += 1
            else:
                for i in range(accomp_min, accomp_max + 1):
                    # formatting or shape should be passed as an argument
                    # if i == 3:
                    #     word_list.append(row[i] + "[~B]") if (row[i] + "[~B]" not in word_list and row[i] != "") else word_list
                    # else:
                    word_list.append(row[i]) if (row[i] not in word_list and row[i] != "") else word_list

                line_count += 1
        
        for word in word_list:
            print(bcolors.OKBLUE + '---- New Shape ----' + bcolors.ENDC)
            print (word)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.margin_left = 0
            tf.margin_right = 0

            p = tf.paragraphs[0]
            run = p.add_run()

            if "|" in word:
                print(word.replace("|", "\n"))
                broken_word = word.split('|')
                run.text = broken_word[0]
                p2 = tf.add_paragraph()
                p2.text = broken_word[1]
                p2_font = p2.font
                p2_font.size = Pt(font_size) #Pt(24)
                p2_font.name = font_name
                p2.level = 1
            else:
                run.text = word

            font = run.font
            print(font_name, font_size)
            font.name = font_name #'Arial'
            font.size = Pt(font_size) #Pt(24)
            if "[~B]" in word:
                run.text = word.replace('[~B]', '')
                font.bold = True
                print("Bold")
                print(run.text)
            # font.bold = True
            # font.italic = None  # cause value to be inherited from theme
            # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            # tf.word_wrap = False
            # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    csv_file.close()

    print(bcolors.OKBLUE + '--------------------' + bcolors.ENDC)
    print(bcolors.HEADER + '\nCreated ' + str(len(word_list)) + ' shapes.\n' + bcolors.ENDC)

    prs.save(ppt_output_filename)

    print(word_list)

    cmd = input('Enter c to continue, x to exit\n')


    if cmd == "c":
        if not os.path.exists('inputs\save.png'):
            print(bcolors.FAIL + 'Save button image file is missing...' + bcolors.ENDC)
            print(bcolors.WARNING + 'Take a screenshot of LibreOffice Impress save button, and save to inputs\save.png' + bcolors.ENDC)
            print('This is screen resolution dependent...')
            sys.exit(1)

        print(bcolors.WARNING + 'Opening ' + ppt_output_filename + ' in libreoffice' + bcolors.ENDC)
        print('When open, let it load, and set focus on the window.')
        print('pyautogui will find the save button and preform a Ctrl + S when "seen"' + bcolors.ENDC)

        args = ['C:\Program Files\LibreOffice\program\soffice', '--impress', '--nologo', ppt_output_filename]
        program = Popen(args, shell=False)

        while True:
            save = pyautogui.locateOnScreen('inputs\save.png')
            if save is not None:
                break
        #program.wait()
        print('Found Save Button')
        #sleep(2)
        hotkey('ctrl', 's')
        sleep(5)
        hotkey('ctrl', 'q')
        # program.terminate()

        prs = Presentation(ppt_output_filename)
        slide = prs.slides[0]
        print(bcolors.WARNING + 'All Shapes:' + bcolors.ENDC)
        list_of_str = []
        with open(output_filename, 'w', encoding='UTF-8', newline='') as write_obj:
            csv_writer = csv.writer(write_obj)
            csv_writer.writerow(["word", "length"])
            line_num = 1
            for shape in slide.shapes:
                list_of_str.append(shape.width)
                csv_writer.writerow([word_list[line_num-1], shape.width])
                line_num += 1

                # Overflow Check
                # shape should be passed as an argument to read this size
                if shape.width >= Inches(7.61):
                    print(bcolors.WARNING + '[W] ' + shape.text + bcolors.ENDC)
            
            print(len(list_of_str))
            print(list_of_str)

        write_obj.close()

        print("Created new csv with word lengths: " + output_filename)

    else:
        sys.exit(1)

if __name__ == "__main__":
    main()