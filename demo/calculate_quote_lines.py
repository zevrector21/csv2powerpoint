from pptx import Presentation
from pptx.util import Inches, Pt
import csv
import os
import sys, getopt
from subprocess import Popen, PIPE, STDOUT
from time import sleep

import pyautogui
from pyautogui import hotkey

def add_column_in_csv(input_file, output_file, transform_row):
    """ Append a column in existing csv using csv.reader / csv.writer classes"""
    with open(input_file, 'r', encoding='UTF-8', newline='') as read_obj, \
            open(output_file, 'w', encoding='UTF-8', newline='') as write_obj:
        csv_reader = csv.reader(read_obj)
        csv_writer = csv.writer(write_obj)
        for row in csv_reader:
            transform_row(row, csv_reader.line_num)
            csv_writer.writerow(row)

def usage():
    print('Usage: calculate_quote_lines.py -i <inputfile> -f <fontname> -s <fontsize>')

def main():
    try:
        opts, args = getopt.getopt(sys.argv[1:],"hi:f:s:",["help","input=","font=","fontsize="])
    except getopt.GetoptError as err:
        print('Unrecognized input')
        print(err)
        usage()
        sys.exit(2)
    csv_filename = None
    ppt_output_filename = None
    font_name = None
    font_size = None
    for o, a in opts:
        if o in ("-h", "--help"):
            usage()
            sys.exit()
        elif o in ("-i", "--input"):
            csv_filename = a
            print("Input: " + csv_filename)
            output_filename = os.path.splitext(csv_filename)[0]+"_quote_heights.csv"
            print("Output: " + output_filename)
            ppt_output_filename = os.path.splitext(csv_filename)[0]+"_quote_test.pptx"
        elif o in ("-f", "--font"):
            font_name = a
            print("Using font: " + font_name)
        elif o in ("-s", "--fontsize"):
            font_size = int(a)
            print("Using font size: " + str(font_size))
        else:
            assert False, "unhandled option"

    if csv_filename == None:
        usage()
        sys.exit()

    os.system('color 0')
    class bcolors:
        HEADER = '\033[95m'
        OKBLUE = '\033[96m'
        OKGREEN = '\033[92m'
        WARNING = '\033[93m'
        FAIL = '\033[91m'
        ENDC = '\033[0m'
        BOLD = '\033[1m'
        UNDERLINE = '\033[4m'

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    word_list = []

    with open(csv_filename, 'r', encoding='UTF-8', newline='') as csv_file:
        csv_reader = csv.reader(csv_file, delimiter=',')
        line_count = 0
        left = top = Inches(0)
        # these should be passed as arguments or pass the shape
        width = Inches(12.34)
        height = Inches(1.85)

        for row in csv_reader:
            if line_count == 0:
                print(bcolors.WARNING + 'Imported CSV Columns:' + bcolors.ENDC)
                print(', '.join(row)) 
                line_count += 1
            else:
                max_length = 0
                for i in range(9, 10):
                    word_list.append(row[i]) if (row[i] not in word_list and row[i] != "") else word_list

                    char_length = len(row[i])

                    if char_length > max_length:
                        max_length_col = i

                line_count += 1
        
        for word in word_list:
            print(bcolors.OKBLUE + '---- New Shape ----' + bcolors.ENDC)

            print (word)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)

            p = tf.paragraphs[0]
            run = p.add_run()

            if "|" in word:
                print(word.replace("|", "\n"))
                p.text = word.replace("|", "\n")
                p_font = p.font
                p_font.size = Pt(font_size) #Pt(24)
                p_font.name = font_name
                p.space_before = 0
                print(p.line_spacing)
                p.line_spacing = 1.2
                # broken_word = word.split('|')
                # run.text = broken_word[0]
                # for word in broken_word:
                #     p2 = tf.add_paragraph()
                #     p2.text = word #broken_word[1]
                #     p2_font = p2.font
                #     p2_font.size = Pt(font_size) #Pt(24)
                #     p2_font.name = font_name
                #     p2.space_before = 0
                #     p2.line_spacing = 1.2

                # p2.level = 1
            else:
                run.text = word

            font = run.font
            print(font_name, font_size)
            font.name = font_name
            font.size = Pt(font_size)

            tf.word_wrap = True

    csv_file.close()

    print(bcolors.OKBLUE + '--------------------' + bcolors.ENDC)
    print(bcolors.HEADER + '\nCreated ' + str(len(word_list)) + ' shapes.\n' + bcolors.ENDC)

    prs.save(ppt_output_filename)

    print(word_list)

    cmd = input('Enter c to continue, x to exit\n')


    if cmd == "c":
        print(bcolors.WARNING + 'Opening ' + ppt_output_filename + ' in libreoffice' + bcolors.ENDC)
        print(bcolors.WARNING + 'Save and close to continue...' + bcolors.ENDC)
        args = ['C:\Program Files\LibreOffice\program\soffice', '--impress', '--nologo', ppt_output_filename]

        program = Popen(args, shell=False)
        while True:
            save = pyautogui.locateOnScreen('save.png')
            if save is not None:
                break
        print('Found Save Button')
        hotkey('ctrl', 's')
        sleep(5)
        hotkey('ctrl', 'q')

        prs = Presentation(ppt_output_filename)
        slide = prs.slides[0]
        print(bcolors.WARNING + 'All Shapes:' + bcolors.ENDC)
        list_of_str = []

        with open(output_filename, 'w', encoding='UTF-8', newline='') as write_obj:
            csv_writer = csv.writer(write_obj)
            csv_writer.writerow(["word", "length"])
            line_num = 1
            for shape in slide.shapes:
                num_lines = round(shape.height/Inches(0.5))
                list_of_str.append(num_lines)
                csv_writer.writerow([word_list[line_num-1], num_lines])
                line_num += 1

                # Overflow Check
                if num_lines >= 5:
                    print(bcolors.FAIL + shape.text + bcolors.ENDC)
                elif num_lines == 4:
                    print(bcolors.WARNING + '[W] ' + shape.text + bcolors.ENDC)

            
            print(len(list_of_str))
            print(list_of_str)

        write_obj.close()

        print("Created new csv with word lengths: " + output_filename)

    else:
        sys.exit(1)

if __name__ == "__main__":
    main()