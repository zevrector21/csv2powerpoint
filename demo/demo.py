import csv
import requests
import os
import sys
import subprocess
import imghdr
import urllib.parse as urlparse
from urllib.parse import parse_qs
import random
import math
import time
from time import sleep
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.dml.color import ColorFormat, RGBColor

from requests.adapters import HTTPAdapter
from requests.packages.urllib3.util.retry import Retry

from pydub import AudioSegment
AudioSegment.ffmpeg = '/ffmpeg/bin/ffmpeg.exe'

from lxml import etree

import png
import pyqrcode

from PIL import Image, ExifTags

# ----------------------------- Setup Vars ---------------------------------
slide_template = 'inputs/Demo Template.pptx'

csv_filename = 'inputs/Demo.csv'
directory_prefix = "inputs"

PLAY_BUTTON_IMG = 'inputs/play_button.png'
PLAY_BUTTON_OFF_SLIDE = True
PLAY_BUTTON_SHADOW = False

# Will automatically convert with whatever is in csv, even both WAV and MP3s mixed
USE_MP3 = False # False = WAV, True = MP3
USE_M4A = True
USE_NOT_PROCESSED = True
CONVERT_M4A = False

# It's nearly impossible to know the width and height of text based on the font size/family/spacing etc.
# Even worse is an auto-fit shape will not auto adjust it's size until opened and clicked in powerpoint (which will then re-render it)
# https://github.com/scanny/python-pptx/issues/147#issuecomment-76804459
# In the XML behind the slide, the current "auto-fitted" font size of the textbox is cached.
# LibreOffice automatically recalculates this cached figure when the presentation opens; PowerPoint does not.
# https://stackoverflow.com/a/41901546
# A workaround is to create a seperate ppt with all names, open it in Libreoffice (which auto-renders on open), measure the lengths, and close it.
# It's possible opening the file in Keynote on a mac will do the same thing
#
# fit_text()
# https://python-pptx.readthedocs.io/en/latest/api/text.html#pptx.text.text.TextFrame.fit_text
# https://python-pptx.readthedocs.io/en/latest/dev/analysis/txt-fit-text.html
# will shrink font size to fit to box with wordwrap, not fit to one line
# It is just as though the user had reduced the font size just until all the text fit within the shape.
#
# auto_size
# SHAPE_TO_FIT_TEXT
# Will resize shape to fit the text, but only on render
# TextFrame.auto_size behaviors of a shape can cause relocation of the shape, in particular, when the "Resize shape to fit text" option is selected,
# one or more of the size or position attributes are changed to comply. 
# https://github.com/scanny/python-pptx/issues/147#issuecomment-76804459
#
MEASURE_NAME_LENGTHS = True # If True, you will be asked to (re)calculate lengths, which will generate a new csv file with a new column containing the lengths
PLAY_NEAR_NAME = False # If True, MEASURE_NAME_LENGTHS will also be True

TITLE_FONT = "Georgia"
TITLE_FONT_SIZE = 40

MEASURE_ACCOMP_LENGTHS = True # If True, you will be asked to (re)calculate lengths, which will generate a new csv file with a new column containing the lengths
SUBTEXT_FONT = "Calibri"
SUBTEXT_FONT_SIZE = 28
SUBTEXT_FONT_SIZE_MINIMUM = 18
SUBTEXT_LINE_SPACING = 0.9
SUBTEXT_SPACE_BEFORE = 10
SUBTEXT_SPACE_AFTER = 0
MOVE_SUBTEXT_UP = False

STUDENT_ID = 0
FULLNAME_COL = 1

HAS_ACCOMPLISHMENTS = True
ACCOMP_TO_TRACK = 2 # Useful for splitting, title slides, etc. Accomp number, not column number.

ACCOMPLISH_1_COL = 2 # Degree
ACCOMPLISH_2_COL = 3 # School
ACCOMPLISH_3_COL = 4 # Latin Honors
ACCOMPLISH_4_COL = 5 # Awards
AUDIO_COL = 6
IMAGE_COL = 7

QUOTE_COL = ""
# QUOTE_FONT = "Calibri" 
# QUOTE_FONT_SIZE = 24
MEASURE_QUOTE_HEIGHT = False

REMOVE_BLANK_IMAGE_PLACEHOLDER = False
MISSING_IMAGE_REPLACEMENT = "inputs/3x4_person_placeholder.png"
USE_FULL_RECT_IMAGE = True
RESIZE_VIDEO = True

STUDENT_SLIDE_LAYOUT_NUM = 1 # Static value or first layout to start iterator at for multi layouts
MULTI_LAYOUT_END_NUM = 2 # Used to iterate or flip layouts starting at STUDENT_SLIDE_LAYOUT_NUM

CREATE_TOC_SLIDE = False
CREATE_TITLE_SLIDES = True
TITLE_SLIDE_LAYOUT_NUM = 0
TOC_SLIDE_LAYOUT_NUM = 3
HAS_TITLE_SLIDE_AUDIO = False
TITLE_SLIDE_AUDIO_FOLDER = ""

AUTO_ADVANCE_SLIDE = True
AUTO_PLAY_AUDIO = True
AUTO_PLAY_VIDEO = True
DEFAULT_SLIDE_DURATION = 4000 # For no audio slides (ms)
SLIDE_GAP = 1.5 # Gap between slides with audio (s)

SPLIT_OUTPUT = False
SPLIT_BY_MASTER = False
SPLIT_BY_NUM = False
SLIDES_PER_PPT = 300

CREATE_QRCODE = False

# ----------------------------- Program Vars ---------------------------------
os.system('color 0')
class bcolors:
    HEADER = '\033[95m'
    OKBLUE = '\033[96m'
    OKGREEN = '\033[92m'
    WARNING = '\033[93m'
    FAIL = '\033[91m'
    ENDC = '\033[0m'
    BOLD = '\033[1m'
    UNDERLINE = '\033[4m'

retry_strategy = Retry(
    total=3,
    status_forcelist=[429, 500, 502, 503, 504],
    method_whitelist=["HEAD", "GET", "OPTIONS"]
)
adapter = HTTPAdapter(max_retries=retry_strategy)
http = requests.Session()
http.mount("https://", adapter)
http.mount("http://", adapter)

# Spoof User-Agent for downloads
http = requests.Session()
http.headers.update({
    "User-Agent": "Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:68.0) Gecko/20100101 Firefox/68.0"
})

is_missing_audio = False
is_missing_image = False
total_missing_photos = 0
total_missing_audio = 0
new_name_shape = None
name_to_length = None
subtext_to_length = None
quote_to_height = None

accomp_min = 99
accomp_max = 0
accomp_col_min = 99
accomp_col_max = 0
accomp_col_range = []

prs_width = 13.333
prs_height = 7.5
px_per_inch = 96
emus_per_inch = 914400

cur_layout = STUDENT_SLIDE_LAYOUT_NUM 

last_track = ''
title_text = ''
title_slide_count = 0
slide_count = 0
title_slides_list = []
tc_slide = None
tc_subtext = None
total_slide_time = 0
cur_master = 0
last_master = ''

start_time = time.perf_counter()
end_time = None
cur_time_start = None
elapsed_time = None
last_time_descp = 'Start'

split_total = 0

warnings_list = []

MY_NAMESPACES={'a' : 'http://schemas.openxmlformats.org/drawingml/2006/main', 'r' : 'http://schemas.openxmlformats.org/officeDocument/2006/relationships', 'p' : 'http://schemas.openxmlformats.org/presentationml/2006/main', 'mc' : 'http://schemas.openxmlformats.org/markup-compatibility/2006', 'p14' : 'http://schemas.microsoft.com/office/powerpoint/2010/main'}

# ----------------------------- Functions ---------------------------------

def get_length(input_audio):
    result = subprocess.run(['ffprobe', '-v', 'error', '-show_entries', 'format=duration', '-of', 'default=noprint_wrappers=1:nokey=1', input_audio], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    return float(result.stdout)

def get_size(input_video):
    result = subprocess.run(['ffprobe', '-v', 'error', '-select_streams', 'v:0', '-show_entries', 'stream=width,height', '-of', 'csv=s=,:p=0', input_video], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    size = result.stdout.decode('utf-8').replace('\r\n','').split(',')
    for i in range(0, len(size)):
        size[i] = int(size[i])
    return size

def create_poster_image(input_video, output_image):
    result = subprocess.run(['ffmpeg', '-i', input_video, '-vf', 'select=eq(n\,0)', '-q:v', '3', output_image], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    return result.returncode

def create_scaled_video(input_video, output_video, width, height):
    result = subprocess.run(['ffmpeg', '-i', input_video, '-vf', 'scale=' + str(width) + ":" +str(height), output_video], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    return result.returncode

def image_variant_exists(input_image):
    image_path = os.path.dirname(input_image) + "/"
    image_name = os.path.splitext(os.path.basename(input_image))[0]

    for extension in ["jpg","jpeg","png","tiff","gif","bmp","mp4","m4v","mov","avi","mpg","mpeg","wmv"]:
        img_variant = image_path + image_name + "." + extension
        if os.path.exists(img_variant):
            return img_variant
    return ''

def add_slide_duration(slide, duration):
    global total_slide_time
    slide_tree = slide._element

    alt_cnt_xml = etree.Element('{%s}AlternateContent' % MY_NAMESPACES['mc'], nsmap=MY_NAMESPACES)
    choice_xml = etree.SubElement(alt_cnt_xml, '{%s}Choice' % MY_NAMESPACES['mc'], nsmap=MY_NAMESPACES)
    
    choice_xml.set("Requires", "p14")
    transition_xml = etree.SubElement(choice_xml, '{%s}transition' % MY_NAMESPACES['p'], nsmap=MY_NAMESPACES)
    transition_xml.set("spd", "slow")
    transition_xml.set("{%s}dur" % MY_NAMESPACES['p14'], "2000")
    transition_xml.set("advClick", "0")
    transition_xml.set("advTm", str(duration))

    fallback_xml = etree.SubElement(alt_cnt_xml, '{%s}Fallback' % MY_NAMESPACES['mc'], nsmap=MY_NAMESPACES)
    transition_f_xml = etree.SubElement(fallback_xml, '{%s}transition' % MY_NAMESPACES['p'], nsmap=MY_NAMESPACES)
    transition_f_xml.set("spd", "slow")
    transition_f_xml.set("advClick", "0")
    transition_f_xml.set("advTm", str(duration))

    slide_tree.insert(2,alt_cnt_xml) # has to go before timing block

    total_slide_time += duration

def add_audio_button(slide, audio_filename, name):
    try:
        icon_size = 0.5
        padding = icon_size * 0.5
        button_width = button_height = Inches(icon_size)

        if PLAY_NEAR_NAME:
            button_left = name.left + name.width
            button_top = name.top + ((name.height-button_height)/2)

            if button_left + button_width >= prs_width*emus_per_inch:
                print(bcolors.WARNING + 'Button out of bounds:' + bcolors.ENDC)
                sys.exit(0)

        else:
            if PLAY_BUTTON_OFF_SLIDE:
                button_left = Inches(-padding - icon_size)
                button_top = Inches(prs_height - icon_size)
            else:
                button_left = Inches(padding)
                button_top = Inches(prs_height - padding - icon_size)
        
        if USE_MP3:
            mtype = 'audio/mp3'
        elif USE_M4A:
            mtype = 'audio/m4a'
        else:
            mtype = 'audio/wav'

        audio = slide.shapes.add_movie(audio_filename, 
            button_left, button_top, button_width, button_height, 
            mime_type = mtype, 
            poster_frame_image = PLAY_BUTTON_IMG)

        if PLAY_BUTTON_SHADOW:
            add_shadow(audio)

        if AUTO_PLAY_AUDIO:
            # Not great if other timings are used...
            tree = audio._element.getparent().getparent().getnext().getnext()
            if movie_filename == '':
                timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
            else:
                timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][1]
            timing.set('delay', '0')

        print('Added Audio Button')

    except Exception as e:
        print(bcolors.FAIL + 'Creating Play Button failed!' + bcolors.ENDC)
        log_error(e)
        print(bcolors.WARNING + 'All Shapes:' + bcolors.ENDC)
        for shape in slide.shapes:
            print('%d %s' % (shape.shape_id, shape.name))
        sys.exit(0)

def add_shadow(shape):
    shadow = shape.shadow
    # This does work, and will generate the blank effectLst element only if set to false
    # https://python-pptx.readthedocs.io/en/latest/api/dml.html#pptx.dml.effect.ShadowFormat
    shadow.inherit = False
    # These don't work:
    # shadow.visible = True
    # shadow.distance = Pt(10)
    # shadow.shadow_type = 'outer'
    # shadow.angle = 45
    # shadow.blur_radius = Pt(5)
    # shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
    # shadow.transparency = '50'
    # shadow.distance = Pt(5)
    # audio.shadow.style = 'outer'
    
    effect_tree = shape._element
    effect = [el for el in effect_tree.iterdescendants() if etree.QName(el).localname == 'effectLst'][0]
    shadow_xml = etree.Element('{%s}outerShdw' % MY_NAMESPACES['a'], nsmap=MY_NAMESPACES)
    shadow_xml.set("blurRad", "50800")
    shadow_xml.set("dist", "38100")
    shadow_xml.set("dir", "2700000")
    shadow_xml.set("algn", "tl")
    shadow_xml.set("rotWithShape", "0")
    shadow_color = etree.SubElement(shadow_xml, '{%s}prstClr' % MY_NAMESPACES['a'], nsmap=MY_NAMESPACES)
    shadow_color.set("val", "black")
    shadow_alpha = etree.SubElement(shadow_color, '{%s}alpha' % MY_NAMESPACES['a'], nsmap=MY_NAMESPACES)
    shadow_alpha.set("val", "40000")
    effect.append( shadow_xml )

def placeholder_error(slide):
    print(bcolors.FAIL + '\nFilling Placeholders failed!' + bcolors.ENDC)
    print(bcolors.WARNING + 'Existing Placeholders:' + bcolors.ENDC)
    for shape in slide.placeholders:
            print('Index: %d Name: %s Type: %s' % (shape.placeholder_format.idx, shape.name, shape.placeholder_format.type))
    sys.exit(0)

def log_error(error):
    print(bcolors.FAIL + '\nERROR LOG' + bcolors.ENDC)
    print(bcolors.FAIL + '----------------------------' + bcolors.ENDC)
    print(bcolors.WARNING + str(error) + bcolors.ENDC)
    exc_type, exc_obj, exc_tb = sys.exc_info()
    fname = os.path.split(exc_tb.tb_frame.f_code.co_filename)[1]
    print(exc_type, fname, exc_tb.tb_lineno)
    print(bcolors.FAIL + '----------------------------\n' + bcolors.ENDC)

def save_presentation(prs, cur):
    outputDir = 'output'
    if not os.path.exists(outputDir):
        os.mkdir(outputDir)
        print(bcolors.HEADER + 'Output Directory Created')
    print(bcolors.ENDC)

    new_ppt_name = os.path.splitext(os.path.split(csv_filename)[1])[0]

    if SPLIT_OUTPUT or cur != "":
        if SPLIT_BY_NUM:
            new_ppt_name = new_ppt_name + "_" + str(cur)
        else:
            new_ppt_name = str(cur).strip()
        new_ppt_name = re.sub(r'[\\/*?:"<>|]',"-",new_ppt_name)
    new_ppt_name = outputDir + "/" + new_ppt_name
    try:
        prs.save(new_ppt_name + '.pptx')
        print(bcolors.WARNING + 'Saving file to:' + bcolors.ENDC)
        print(new_ppt_name + '.pptx\n')
    except Exception as e:
        print(bcolors.FAIL + 'Failed to save file.' + bcolors.ENDC + ' Is it currently open in PPT?')
        log_error(e)
        increm = 0
        incremental_filename = new_ppt_name + '.pptx'
        while os.path.exists(incremental_filename):
            increm += 1
            incremental_filename = new_ppt_name + '_' + str(increm) + '.pptx'
        print(bcolors.WARNING + 'Saving file to:' + bcolors.ENDC)
        print(incremental_filename)
        prs.save(incremental_filename)

def csv_sanity_check(csv_file):
    if not os.path.exists(csv_file):
        print(bcolors.FAIL + 'CSV File: "' + csv_file + '" not found!' + bcolors.ENDC)
        sys.exit(0)
    else:
        print(bcolors.OKBLUE + "Using CSV: " + csv_file + bcolors.ENDC)

def template_sanity_check(template_file):
    if not os.path.exists(template_file):
        print(bcolors.FAIL + 'Powerpoint Template File: "' + template_file + '" not found!' + bcolors.ENDC)
        sys.exit(0)
    else:
        print(bcolors.OKBLUE + "Using PPT Template: " + template_file + bcolors.ENDC)

def fix_image_orientation(image_file):
    image=Image.open(image_file)
    rotation = 0
    try:
        for orientation in ExifTags.TAGS.keys():
            if ExifTags.TAGS[orientation]=='Orientation':
                break
        exif=dict(image._getexif().items())

        # 1: 'Horizontal (normal)',
        # 2: 'Mirrored horizontal',
        # 3: 'Rotated 180',
        # 4: 'Mirrored vertical',
        # 5: 'Mirrored horizontal then rotated 90 CCW',
        # 6: 'Rotated 90 CW',
        # 7: 'Mirrored horizontal then rotated 90 CW',
        # 8: 'Rotated 90 CCW'

        if exif[orientation] == 3:
            image=image.transpose(Image.ROTATE_180)
            rotation = 180
        elif exif[orientation] == 6:
            image=image.transpose(Image.ROTATE_270)
            rotation = 90
        elif exif[orientation] == 8:
            image=image.transpose(Image.ROTATE_90)
            rotation = 270
        else:
            rotation = 0
        image.save(image_file)
    except (AttributeError, KeyError, IndexError):
        pass
    image.close()
    return rotation

def fix_video_orientation(video_file):
    rotation = -1
    try:
        result = subprocess.run(['ffprobe', '-v', 'error', '-select_streams', 'v:0', '-show_entries', 'stream_tags=rotate', '-of', 'default=nw=1:nk=1', '-i', video_file], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
        rotation = result.stdout.decode('utf-8').replace('\r\n','')

        if rotation != 0 and rotation != '':
            video_path = os.path.dirname(video_file) + "/"
            video_name = os.path.splitext(os.path.basename(video_file))[0]
            extension = os.path.splitext(video_file)[1]
            video_variant = video_path + video_name + "_rot" + extension
            result = subprocess.run(['ffmpeg', '-i', video_file, '-c:a', 'copy', video_variant], stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
            os.remove(video_file)
            os.rename(video_variant, video_file)
        rotation = 0
    except (AttributeError, KeyError, IndexError):
        pass
    return rotation
    
def call_time(descp):
    descp = str(descp)
    global start_time, cur_time_start, elapsed_time, last_time_descp
    if cur_time_start == None:
        elapsed_time = time.perf_counter() - start_time
    else:
        elapsed_time = time.perf_counter() - cur_time_start
    
    print(bcolors.OKGREEN + "Time Elapsed: " + str(elapsed_time) + ' from: ' + last_time_descp + ' to: ' + descp)

    cur_time_start = time.perf_counter()
    last_time_descp = descp
    return elapsed_time

def millis_to_hms(millis):
    millis = int(millis)
    seconds=(millis/1000)%60
    seconds = int(seconds)
    minutes=(millis/(1000*60))%60
    minutes = int(minutes)
    hours=(millis/(1000*60*60))%24
    return("%d:%d:%d\n" % (hours, minutes, seconds))

def load_accomp_range():
    global accomp_min, accomp_max, accomp_col_min, accomp_col_max, accomp_col_range
    accomp_test = 1
    cur_col = 0
    while True:
        if 'ACCOMPLISH_' + str(accomp_test) + '_COL' in globals():
            if accomp_test < accomp_min:
                accomp_min = accomp_test
            if accomp_test > accomp_max:
                accomp_max = accomp_test
            
            cur_col = globals()['ACCOMPLISH_' + str(accomp_max) + '_COL']
            if cur_col < accomp_col_min:
                accomp_col_min = cur_col
            if cur_col > accomp_col_max:
                accomp_col_max = cur_col

            accomp_test += 1
        else:
            break  
    accomp_col_range = str(accomp_col_min) + "," + str(accomp_col_max)
    return accomp_col_range

# ---------------------------------------------------------------------------

os.system('cls' if os.name == 'nt' else 'clear') # clear terminal

print(bcolors.OKGREEN)
print(" ___ ___ _______  __   ___                       _           ")
print("| _ \ _ \_   _\ \/ /  / __|___ _ _  ___ _ _ __ _| |_ ___ _ _ ")
print("|  _/  _/ | |  >  <  | (_ / -_) ' \/ -_) '_/ _` |  _/ _ \ '_|")
print("|_| |_|   |_| /_/\_\  \___\___|_||_\___|_| \__,_|\__\___/_|  ")
print("Written by Jon Olsen\n")
print(bcolors.ENDC)

print(bcolors.OKBLUE + "Current working directory: " + os.getcwd() + bcolors.ENDC)
csv_sanity_check(csv_filename)
template_sanity_check(slide_template)

print("\n")

if PLAY_NEAR_NAME:
    MEASURE_NAME_LENGTHS = True
if MEASURE_NAME_LENGTHS:
    while True:
        temp_csv_filename = os.path.splitext(csv_filename)[0]+"_name_lengths.csv"
        if not os.path.exists(temp_csv_filename):
            print(bcolors.WARNING + 'Name lengths not calculated yet!' + bcolors.ENDC)
            cmd = input('Run calculation program? (y/n)\n')
            if cmd == "n":
                sys.exit(0)
        else:
            cmd = input('Name lengths file found, do you want to recalculate? (y/n/q)\n')
            if cmd == "n":
                with open(temp_csv_filename, 'r', encoding='UTF-8', newline='') as f:
                    reader = csv.DictReader(f)
                    name_to_length = {row['word']: row["length"] for row in reader}   
                f.close()
                print("Loaded name length dictionary...") 
                break
            if cmd == "q":
                sys.exit(0)
        print(bcolors.OKGREEN + 'Running name width calculation...' + bcolors.ENDC)
        subprocess.call(['py', 'calculate_name_lengths.py', '-i' + csv_filename, '-c' + str(FULLNAME_COL), '-f' + TITLE_FONT, '-s' + str(TITLE_FONT_SIZE)])
load_accomp_range()
if MEASURE_ACCOMP_LENGTHS: 
    while True:
        temp_csv_filename = os.path.splitext(csv_filename)[0]+"_accomp_lengths.csv"
        if not os.path.exists(temp_csv_filename):
            print(bcolors.WARNING + 'Accomplishment lengths not calculated yet!' + bcolors.ENDC)
            cmd = input('Run calculation program? (y/n)\n')
            if cmd == "n":
                sys.exit(0)
        else:
            cmd = input('Accomplishment lengths file found, do you want to recalculate? (y/n/q)\n')
            if cmd == "n":
                with open(temp_csv_filename, 'r', encoding='UTF-8', newline='') as f:
                    reader = csv.DictReader(f)
                    subtext_to_length = {row['word']: row["length"] for row in reader}
                f.close()
                print("Loaded word length dictionary...")
                break
            if cmd == "q":
                sys.exit(0)
        print(bcolors.OKGREEN + 'Running accomplishment width calculation...' + bcolors.ENDC)
        print('py', 'calculate_accomp_lengths.py', '-i' + csv_filename, '-f' + SUBTEXT_FONT, '-s' + str(SUBTEXT_FONT_SIZE))
        subprocess.call(['py', 'calculate_accomp_lengths.py', '-i' + csv_filename, '-r' + accomp_col_range, '-f' + SUBTEXT_FONT, '-s' + str(SUBTEXT_FONT_SIZE)])
if MEASURE_QUOTE_HEIGHT:
    while True:
        temp_csv_filename = os.path.splitext(csv_filename)[0]+"_quote_heights.csv"
        if not os.path.exists(temp_csv_filename):
            print(bcolors.WARNING + 'Quote heights not calculated yet!' + bcolors.ENDC)
            cmd = input('Run calculation program? (y/n)\n')
            if cmd == "n":
                sys.exit(0)
        else:
            cmd = input('Quote heights file found, do you want to recalculate? (y/n/q)\n')
            if cmd == "n":
                # csv_filename = temp_csv_filename
                with open(temp_csv_filename, 'r', encoding='UTF-8', newline='') as f:
                    reader = csv.DictReader(f)
                    quote_to_height = {row['word']: row["length"] for row in reader}
                f.close()
                print("Loaded word length dictionary...")
                break
            if cmd == "q":
                sys.exit(0)
        print(bcolors.OKGREEN + 'Running quote height calculation...' + bcolors.ENDC)
        # subprocess.call(['py', 'calculate_accomp_lengths.py', '-i' + csv_filename])
        print('py', 'calculate_quote_lines.py', '-i' + csv_filename, '-f' + QUOTE_FONT, '-s' + str(QUOTE_FONT_SIZE))
        subprocess.call(['py', 'calculate_quote_lines.py', '-i' + csv_filename, '-f' + QUOTE_FONT, '-s' + str(QUOTE_FONT_SIZE)])

prs = Presentation(slide_template)
prs_height = prs.slide_height/emus_per_inch
prs_width = prs.slide_width/emus_per_inch
px_per_inch = round(1920 / prs_width)

with open(csv_filename, 'r', encoding='UTF-8', newline='') as csv_file:
    print(bcolors.WARNING + 'Opening: ' + csv_filename + bcolors.ENDC)
    csv_reader = csv.reader(csv_file, delimiter=',')
    line_count = 0
    files_downloaded = False # Flag used for delaying to limit server hit rate 
    for row in csv_reader:
        if line_count == 0:
            print(bcolors.WARNING + 'Imported CSV Columns:' + bcolors.ENDC)
            print(', '.join(row)) 
            line_count += 1

            # Create prefix directory
            if not os.path.exists(directory_prefix):
                os.mkdir(directory_prefix)
                print(bcolors.HEADER + 'Prefix Directory Created')

            # Create Audio Recordings Directory if don't exist
            audioDir = directory_prefix + '/recordings'
            wavDir = audioDir + '/wav/'
            m4aDir = audioDir + '/m4a/'
            upDir = audioDir + '/unprocessed/'
            procDir = upDir + "processed/"

            if not os.path.exists(audioDir):
                os.mkdir(audioDir)
                print(bcolors.HEADER + 'MP3 Audio Directory Created')

            # Create Images Directory
            imageDir = directory_prefix + '/images'
            if not os.path.exists(imageDir):
                os.mkdir(imageDir)
                print(bcolors.HEADER + 'Image Directory Created')
            print(bcolors.ENDC)

            # Create Table of Contents Slide
            if CREATE_TOC_SLIDE and CREATE_TITLE_SLIDES and not SPLIT_OUTPUT:
                try:
                    toc_slide_layout = prs.slide_layouts[TOC_SLIDE_LAYOUT_NUM]
                    tc_slide = prs.slides.add_slide(toc_slide_layout)
                    tc_slide.name = "toc_" + str(title_slide_count)
                    
                    tc_title = tc_slide.shapes.title
                    tc_title.text = "Table of Contents"
                    tc_title.top = Inches(0.5)
                    tc_title.width = int(Inches(prs_width)*0.8)
                    tc_title.height = Pt(60)
                    tc_title.left = int((Inches(prs_width)*0.2)/2)
                    tc_subtext = tc_slide.placeholders[1]
                    tc_subtext._element.getparent().remove(tc_subtext._element)

                    title_slide_count += 1
                    slide_count += 1

                    title_slides_list.append(slide_count)

                    if AUTO_ADVANCE_SLIDE:
                        add_slide_duration(tc_slide, DEFAULT_SLIDE_DURATION)

                except Exception as e:
                    print(bcolors.WARNING + 'Table of Contents Placeholder Error...' + bcolors.ENDC)
                    log_error(e)
                    placeholder_error(tc_slide)

        else:
            call_time('Student Row Start')
            if SPLIT_OUTPUT and SPLIT_BY_NUM and ((slide_count-1) % SLIDES_PER_PPT == 0):
                print(bcolors.OKBLUE + 'Splitting Presentation: ' + str(split_total) + bcolors.ENDC)
                save_presentation(prs, "[" + str(split_total*SLIDES_PER_PPT+1) + "-" + str(slide_count-1) + "]")
                split_total += 1
                prs = Presentation(slide_template)

            print(bcolors.OKBLUE + '---- Slide ' + str(slide_count) + ' ----' + bcolors.ENDC)

            print(str(row[FULLNAME_COL]))

            filename = row[STUDENT_ID]

            # ----------------------------------------------------------------------------------------
            # Check Audio
            # call_time('Audio Check')
            audio_url = row[AUDIO_COL]
            if audio_url != '':
                audio_path = urlparse.urlparse(audio_url).path
                extension = os.path.splitext(audio_path)[1]
                
                if extension == '.mp3':
                    if USE_NOT_PROCESSED and "admin_recording" in audio_url:
                        if not os.path.exists(upDir):
                            os.mkdir(upDir)
                            os.mkdir(procDir)
                            print(bcolors.HEADER + 'Unprocessed Audio Directory Created' + bcolors.ENDC)
                        audio_filename = os.path.join(procDir, filename + '.m4a')
                        print(bcolors.HEADER + 'Checking if Processed Audio Exists: ' + audio_filename + bcolors.ENDC)
                        if not os.path.exists(audio_filename):
                            audio_filename = os.path.join(upDir, filename + extension)
                            print(bcolors.WARNING + 'Using Unprocessed Audio...' + bcolors.ENDC)
                            warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Using Unprocessed Audio")
                    else:
                        audio_filename = os.path.join(audioDir, filename + extension)
                elif extension == '.wav':
                    # Create WAV Directory
                    if not os.path.exists(wavDir):
                        os.mkdir(wavDir)
                        print(bcolors.HEADER + 'WAV Audio Directory Created' + bcolors.ENDC)
                    audio_filename = os.path.join(wavDir, filename + extension)

                    # Fix NC wrong filetype for processed wav files
                    if USE_M4A:
                        extension = '.m4a'
                        if not os.path.exists(m4aDir):
                            os.mkdir(m4aDir)
                            print(bcolors.HEADER + 'M4A Audio Directory Created' + bcolors.ENDC)
                        audio_filename = os.path.join(m4aDir, filename + extension)
                    
                    # Use unprocessed audio if exists
                    if USE_NOT_PROCESSED and "normalized_audio" not in audio_url:
                        if not os.path.exists(upDir):
                            os.mkdir(upDir)
                            os.mkdir(procDir)
                            print(bcolors.HEADER + 'Unprocessed Audio Directory Created' + bcolors.ENDC)
                        audio_filename = os.path.join(procDir, filename + extension)
                        print(bcolors.HEADER + 'Checking if Processed Audio Exists: ' + audio_filename + bcolors.ENDC)
                        if not os.path.exists(audio_filename):
                            audio_filename = os.path.join(upDir, filename + extension)
                            print(bcolors.WARNING + 'Using Unprocessed Audio...' + bcolors.ENDC)
                            warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Using Unprocessed Audio")
                else:
                    print(bcolors.FAIL + 'Conversion not handled for: "' + extension + '" audio types' + bcolors.ENDC)
                    print('For', row[FULLNAME_COL], 'Row:', line_count)
                    sys.exit(0)

                # Download Audio if it doesn't exist
                if not os.path.exists(audio_filename):
                    print(bcolors.WARNING + 'Downloading Audio...' + bcolors.ENDC)
                    warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Downloaded New Audio")

                    try:
                        r = http.get(audio_url)
                    except Exception as e:
                        print(bcolors.FAIL + 'Audio download failed!' + bcolors.ENDC)
                        print('For', row[FULLNAME_COL], 'Row:', line_count)
                        log_error(e)
                        sys.exit(0)
                    
                    print(bcolors.OKGREEN + 'Downloaded: ' + bcolors.ENDC + audio_filename)

                    with open(audio_filename, 'wb') as f:
                        f.write(r.content)
                        f.close()

                    files_downloaded = True
                if not USE_M4A and extension != '.m4a':
                    # Convert Audio if needed
                    if extension != '.mp3' and USE_MP3:
                        # Convert WAV to MP3
                        mp3_filename = audioDir + '/' + filename + '.mp3'
                        if not os.path.exists(mp3_filename):
                            print(bcolors.WARNING + 'Converting to MP3...' + bcolors.ENDC)
                            sound = AudioSegment.from_wav(audio_filename)
                            sound.export(mp3_filename, format='mp3')
                            audio_filename = mp3_filename
                            print(bcolors.OKGREEN + 'Converted: ' + bcolors.ENDC + audio_filename)
                    elif extension != '.wav' and not USE_MP3:
                        # Convert MP3 to WAV
                        wav_filename = wavDir + '/' + filename + '.wav'
                        if not os.path.exists(wav_filename):
                            print(bcolors.WARNING + 'Converting to WAV...' + bcolors.ENDC)
                            sound = AudioSegment.from_mp3(audio_filename)
                            sound.export(wav_filename, format='wav')
                            audio_filename = wav_filename
                            print(bcolors.OKGREEN + 'Converted: ' + bcolors.ENDC + audio_filename)
                elif CONVERT_M4A:
                    wav_filename = wavDir + '/' + filename + '.wav'
                    if not os.path.exists(wav_filename):
                        print(bcolors.WARNING + 'Converting M4A to WAV...' + bcolors.ENDC)
                        sound = AudioSegment.from_m4a(audio_filename)
                        sound.export(wav_filename, format='wav')
                        audio_filename = wav_filename
                        print(bcolors.OKGREEN + 'Converted: ' + bcolors.ENDC + audio_filename)
                else:
                    print(bcolors.WARNING + 'Using m4a audio...' + bcolors.ENDC + audio_filename)
            else:
                print(bcolors.FAIL + 'Student audio is missing...' + bcolors.ENDC)
                total_missing_audio += 1
                is_missing_audio = True
                warnings_list.append(str(slide_count) + ": Missing Audio - " + row[FULLNAME_COL])
            
            # ----------------------------------------------------------------------------------------           
            # Download Images
            # call_time('Download Images')
            image_basename = os.path.basename(row[IMAGE_COL])
            image_hash = urlparse.unquote(os.path.splitext(image_basename)[0])
            image_filename = os.path.join(imageDir, filename + '#' + image_hash + '.jpeg')
            movie_filename = ''
            
            # Legacy Image File Name Check
            legacy_image_filename = os.path.join(imageDir, filename + '.jpeg')
            if os.path.exists(legacy_image_filename):
                if os.path.exists(image_filename):
                    os.rename(image_filename, os.path.join(imageDir, filename + '_old.jpeg'))
                os.rename(legacy_image_filename, image_filename)

            if not os.path.exists(image_filename):
                print(bcolors.WARNING + "Cropped JPEG Image Doesn't Exist...")
                image_url = row[IMAGE_COL]
                if image_url != '' and image_url != 'https://www.name-coach.com/images/photo-blank.gif':
                    image_path = urlparse.urlparse(image_url).path
                    extension = os.path.splitext(image_path)[1]

                    if extension == '':
                        if 'drive.google.com' in image_url:
                            image_hash = parse_qs(urlparse.urlparse(image_url).query)['id'][0]
                            extension = ".jpeg"
                        else:
                            print('For', row[FULLNAME_COL], "Can't parse extension: ", image_url)
                            sys.exit(0)

                    temp_dir = imageDir + '/new'
                    if not os.path.exists(temp_dir):
                        os.mkdir(temp_dir)
                        print(bcolors.HEADER + 'Temp Image Directory Created')

                    image_filename = os.path.join(temp_dir, filename + '#' + image_hash + extension)

                    # Legacy Uncropped Image Check
                    legacy_image_filename = os.path.join(temp_dir, filename + extension)
                    if os.path.exists(legacy_image_filename):
                        os.rename(legacy_image_filename, image_filename)

                    # Uncropped Image Check
                    img_variant = image_variant_exists(image_filename)
                    if img_variant == '':
                        print(bcolors.WARNING + 'Downloading New Image...')

                        r = http.get(image_url)

                        if r.status_code == 200:

                            with open(image_filename, 'wb') as f:
                                f.write(r.content)
                                f.close()

                            real_extension = imghdr.what(image_filename)

                            if real_extension != extension.replace('.','') and real_extension != None:
                                fixed_filename = os.path.join(temp_dir, filename + '#' + image_hash.replace(extension,'') + "." + real_extension)
                                os.rename(image_filename, fixed_filename)
                                # if not os.path.exists(fixed_filename):
                                #     os.rename(image_filename, fixed_filename)
                                # else:
                                #     os.remove(image_filename)
                                image_filename = fixed_filename

                            if extension.replace('.','').lower() in ["jpg","jpeg","png","tiff","gif","bmp"]:
                                fix_image_orientation(image_filename)
                            elif extension.replace('.','').lower() in ["mp4","m4v","mov","avi","mpg","mpeg","wmv"]:
                                # ["mp4","m4v","mov","avi","mpg","mpeg","wmv"]
                                fix_video_orientation(image_filename)
                            else:
                                warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Unsupported File: " + image_filename)
                            files_downloaded = True
                            print(bcolors.OKGREEN + 'Downloaded: ' + bcolors.ENDC + image_filename)
                            warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Downloaded New Image")

                        else:
                            print(bcolors.FAIL + 'Image download failed!' + bcolors.ENDC)
                            print('For', row[FULLNAME_COL], 'Row:', line_count)
                            warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Image Download Failed, Status: " + str(r.status_code))
                            is_missing_image = True

                    
                    else:
                        image_filename = img_variant
                        print(bcolors.WARNING + 'Using Temp Image...' + bcolors.ENDC)
                        warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Using Temp Image")
            
                else:
                    print(bcolors.FAIL + 'Image is missing or generic in CSV...' + bcolors.ENDC)
                    total_missing_photos += 1
                    if not REMOVE_BLANK_IMAGE_PLACEHOLDER and MISSING_IMAGE_REPLACEMENT != "":
                        image_filename = MISSING_IMAGE_REPLACEMENT
                    else:
                        is_missing_image = True
            
            # ----------------------------------------------------------------------------------------
            # Create a title slide
            # call_time('Title Slides')
            if CREATE_TITLE_SLIDES:
                if last_track != row[locals()['ACCOMPLISH_' + str(ACCOMP_TO_TRACK) + '_COL']]:
                    if SPLIT_OUTPUT and SPLIT_BY_MASTER:
                        print(bcolors.OKBLUE + 'Splitting Presentation: ' + str(split_total) + bcolors.ENDC)
                        save_presentation(prs, "[" + str(last_track) + "] " + title_text)
                        split_total += 1
                        prs = Presentation(slide_template) 
                    try:
                        last_track = row[locals()['ACCOMPLISH_' + str(ACCOMP_TO_TRACK) + '_COL']]

                        title_slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT_NUM]

                        title_slide = prs.slides.add_slide(title_slide_layout)
                        title_slide.name = "title_" + str(title_slide_count)
                        
                        # ----------------------------------------------------------------
                        # Match Placeholders
                        title_text = last_track # add other accomps or text here
                        title = title_slide.shapes.title
                        # title.text = title_text
                        title_tf = title.text_frame
                        title_p = title_tf.paragraphs[0]
                        title_run = title_p.add_run()
                        title_run.text = title_text

                        if HAS_TITLE_SLIDE_AUDIO:
                            ts_audio_file = TITLE_SLIDE_AUDIO_FOLDER + last_track + ".wav"
                            if os.path.exists(ts_audio_file):
                                add_audio_button(title_slide, ts_audio_file, title)
                                if AUTO_ADVANCE_SLIDE:
                                    ts_audio_dur = get_length(ts_audio_file)
                                    ts_slide_dur = int((ts_audio_dur + SLIDE_GAP)*1000)
                                    print('Adding slide duration: ' + str(ts_slide_dur/1000))
                                    add_slide_duration(title_slide, ts_slide_dur)
                            else:
                                print(bcolors.FAIL + 'Failed to add title slide audio:' + bcolors.ENDC)
                                print('For', ts_audio_file)
                                add_slide_duration(title_slide, DEFAULT_SLIDE_DURATION)
                                warnings_list.append(ts_audio_file + " - Title Slide Audio Missing")
                        elif AUTO_ADVANCE_SLIDE:
                            add_slide_duration(title_slide, DEFAULT_SLIDE_DURATION)
                        
                        title_slide_count += 1
                        slide_count += 1

                        title_slides_list.append(str(slide_count) + ":" + str(cur_master))

                        if CREATE_TOC_SLIDE and not SPLIT_OUTPUT:
                            left = (Inches(prs_width)*.2)/2
                            top = tc_slide.shapes.title.top + tc_slide.shapes.title.height - Pt(14) + (Pt(28) * (title_slide_count - 1)) # * 1.5)
                            width = (Inches(prs_width)*0.8)
                            height = Pt(28)

                            new_link = tc_slide.shapes.add_textbox(left, top, width, height)
                            new_link.click_action.target_slide = title_slide
                            new_link_tf = new_link.text_frame
                            new_link_p = new_link_tf.paragraphs[0]
                            run = new_link_p.add_run()
                            run.text = title_text
                            font = run.font
                            font.size = Pt(14)
                            font.underline = True
                            font.color.rgb = RGBColor(255, 255, 255)
                            new_link_p.alignment = PP_ALIGN.CENTER
                            new_link_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        
                        print(bcolors.OKBLUE + '---- Slide ' + str(title_slide_count) + ' ----' + bcolors.ENDC)
                        print(title_text)
                    except Exception as e:
                        print(bcolors.WARNING + 'Title Slide Placeholder Error...' + bcolors.ENDC)
                        log_error(e)
                        placeholder_error(title_slide)
            
            # ----------------------------------------------------------------------------------------
            # Create a QR Code
            if CREATE_QRCODE:
                qr_dir = imageDir + '/qr'
                qr_path = qr_dir + '/' + row[STUDENT_ID] + ".png"
                if not os.path.exists(qr_dir):
                    os.mkdir(qr_dir)
                    print(bcolors.HEADER + 'QR Code Directory Created' + bcolors.ENDC)
                
                if not os.path.exists(qr_path):
                    qr_code = pyqrcode.create(row[STUDENT_ID], error='M', mode='binary')
                    qr_code.png(qr_path, scale=6, quiet_zone=1)
                    print('Made QR code' + bcolors.ENDC)
            
            # ----------------------------------------------------------------------------------------
            # Create new slide and fill placholders
            # call_time('Fill Placeholder Start')
            try:
                slide_layout = prs.slide_masters[cur_master].slide_layouts[cur_layout]
                slide = prs.slides.add_slide(slide_layout)

                slide.name = row[STUDENT_ID]
              
                # *** Match Placeholders Here ***
                name = slide.placeholders[0]
                subtext = slide.placeholders[1]
                # if row[QUOTE_COL] != "":
                #     quote = slide.placeholders[13]

                # if not is_missing_image:
                #     picture = slide.placeholders[10]
                picture = slide.placeholders[10]

                # call_time('Name Placeholder')
                if MEASURE_NAME_LENGTHS:
                    name_width = int(name_to_length.get(row[FULLNAME_COL], 0))
                    if name_width == 0:
                        print(bcolors.WARNING + 'Name Length Missing' + bcolors.ENDC)
                        print(row[FULLNAME_COL])
                        print(bcolors.FAIL + 'Recalculate Name Lengths...' + bcolors.ENDC)
                        sys.exit(0)
                    if float(name_width) + name.text_frame.margin_left + name.text_frame.margin_right + Inches(0.1) >= name.width:
                        print(bcolors.WARNING + 'Name out of bounds...' + bcolors.ENDC)

                        name.text_frame.text = row[FULLNAME_COL]
                        font = name.text_frame.paragraphs[0].runs[0].font
                        new_font_size = math.floor(((name.width-(name.text_frame.margin_left+name.text_frame.margin_right))/int(name_width))*TITLE_FONT_SIZE)-1

                        print(bcolors.WARNING + str(new_font_size) + bcolors.ENDC)
                        if new_font_size > TITLE_FONT_SIZE :
                            new_font_size = TITLE_FONT_SIZE
                        # else:
                        #     warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Shrinking Name - Size:" + str(new_font_size) + "pt")

                        font.size = Pt(new_font_size)
                    else:
                        name.text = row[FULLNAME_COL]
                else:
                    name.text = row[FULLNAME_COL]

                # call_time('ACCOMP Placeholder')
                if HAS_ACCOMPLISHMENTS:
                    accomp_lines = []
                    for accomp in range(accomp_min, accomp_max + 1):
                        accomp_lines.append(row[locals()['ACCOMPLISH_' + str(accomp) + '_COL']])
                    accomplishments = list(filter(None, accomp_lines))

                    if MEASURE_ACCOMP_LENGTHS and len(accomplishments) > 0:
                        max_accomp_length = 0
                        max_accomp_text = ""
                        accomps = 0

                        for accomp in accomplishments:
                            accomp_length = int(subtext_to_length.get(accomp, 0))
                            if accomp_length == 0:
                                print(bcolors.WARNING + 'Accomplishment Length Missing' + bcolors.ENDC)
                                print(accomp)
                                print(bcolors.FAIL + 'Recalculate Accomplishment Lengths...' + bcolors.ENDC)
                                sys.exit(0)
                            if accomp_length > max_accomp_length:
                                max_accomp_length = accomp_length
                                max_accomp_text = accomp

                            accomps += 1

                        if (max_accomp_length + subtext.text_frame.margin_left + subtext.text_frame.margin_right) >= subtext.width - Inches(0.5):

                            print(bcolors.WARNING + 'Accomplishment out of bounds...' + bcolors.ENDC)

                            # new_font_size = math.floor(((subtext.width-(subtext.text_frame.margin_left+subtext.text_frame.margin_right))/int(max_accomp_length))*SUBTEXT_FONT_SIZE)
                            new_font_size = math.floor(((subtext.width-(subtext.text_frame.margin_left+subtext.text_frame.margin_right))/int(max_accomp_length))*SUBTEXT_FONT_SIZE) - 0.2

                            if (new_font_size < SUBTEXT_FONT_SIZE_MINIMUM):
                                print(bcolors.FAIL + 'Font Size Too Small: ' + str(new_font_size) + bcolors.ENDC)
                                new_font_size = SUBTEXT_FONT_SIZE_MINIMUM
                                warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Accomplishment Minimum Font Size, Potential Edge Overflow")
                            
                            if (new_font_size > SUBTEXT_FONT_SIZE):
                                new_font_size = SUBTEXT_FONT_SIZE

                            warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Shrinking Accomplishment")
                        else:
                            new_font_size = SUBTEXT_FONT_SIZE

                        subtext_p = subtext.text_frame.paragraphs[0]

                        lines_needed = 0

                        if '|' in accomplishments[0]:
                            print(bcolors.WARNING + 'Breaking Accomplishment onto 2 lines' + bcolors.ENDC)
                            broken_word = accomplishments[0].split('|')
                            subtext_p.text = broken_word[0]
                            subtext_p2 = subtext.text_frame.add_paragraph()
                            subtext_p2.text = broken_word[1]
                            subtext_p2_font = subtext_p2.font
                            subtext_p2_font.size = Pt(new_font_size)
                            subtext_p2.level = 1
                            lines_needed += 2
                        else:
                            subtext_p.text = accomplishments[0]
                            lines_needed += 1
                        
                        subtext_p_font = subtext_p.font
                        subtext_p_font.size = Pt(new_font_size)
                        subtext_p.line_spacing = SUBTEXT_LINE_SPACING
                        subtext_p.space_before = Pt(SUBTEXT_SPACE_BEFORE)
                        subtext_p.space_after = Pt(SUBTEXT_SPACE_AFTER)

                        print(accomplishments[0])

                        for para_str in accomplishments[1:]:
                            
                            p = subtext.text_frame.add_paragraph()

                            if '|' in para_str:
                                print(bcolors.WARNING + 'Breaking Accomplishment onto 2 lines' + bcolors.ENDC)
                                broken_word = para_str.split('|')
                                p.text = broken_word[0]
                                p2 = subtext.text_frame.add_paragraph()
                                p2.text = broken_word[1]
                                p2_font = p2.font
                                p2_font.size = Pt(new_font_size)
                                p2.level = 1
                                lines_needed += 2
                            else:
                                p.text = para_str
                                lines_needed += 1
                            
                            p_font = p.font
                            p_font.size = Pt(new_font_size)
                            p.line_spacing = SUBTEXT_LINE_SPACING
                            p.space_before = Pt(SUBTEXT_SPACE_BEFORE)
                            p.space_after = Pt(SUBTEXT_SPACE_AFTER)
                            # p_font.bold = True
                            # p_font.italic = None
                            # p_font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                            # p.level = 1

                            print(para_str)

                        print(bcolors.HEADER + 'Subtext Lines needed: ' + str(lines_needed) + bcolors.ENDC)

                        # -------------------------------------------------------------------

                        if MOVE_SUBTEXT_UP:
                            if lines_needed > 3:
                                print(bcolors.WARNING + 'Moving Title and Subtext Up' + bcolors.ENDC)
                                sub_top = subtext.top
                                sub_left = subtext.left
                                sub_width = subtext.width
                                sub_height = subtext.height
                                sub_line_height = subtext.text_frame.paragraphs[0].line_spacing

                                name_top = name.top
                                name_left = name.left
                                name_width = name.width
                                name_height = name.height

                                sub_offset = ((lines_needed - 3) * Pt(SUBTEXT_FONT_SIZE * sub_line_height))

                                subtext.top = sub_top - sub_offset
                                subtext.left = sub_left
                                subtext.width = sub_width
                                subtext.height = sub_height + sub_offset

                                name.top = name_top - sub_offset
                                name.left = name_left
                                name.width = name_width
                                name.height = name_height

                                if (name.top < Pt(SUBTEXT_FONT_SIZE)):
                                    print(bcolors.FAIL + 'Too close to edge!' + bcolors.ENDC)
                                    # sys.exit(1)
                                    warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Too Close to Edge")
                                
                    else:
                        subtext.text = '\n'.join(filter(None, accomplishments))

                # call_time('Image Placeholder')
                if is_missing_image:
                    if REMOVE_BLANK_IMAGE_PLACEHOLDER:
                        print(bcolors.WARNING + 'Removing image placeholder from slide...' + bcolors.ENDC)
                        picture._element.getparent().remove(picture._element)
                    elif MISSING_IMAGE_REPLACEMENT != "":
                        picture.insert_picture(image_filename)
                    else:
                        print(bcolors.WARNING + 'Missing image and no replacement provided!' + bcolors.ENDC)
                else:
                    media_extension = os.path.splitext(os.path.basename(image_filename))[1].replace('.','').lower()
                    media_name = os.path.splitext(image_filename)[0]
                    try:
                        if USE_FULL_RECT_IMAGE:

                            picture_left = picture.left
                            picture_top = picture.top

                            if media_extension in ["jpg","jpeg","png","tiff","gif","bmp"]:
                                picture_ph = picture.insert_picture(image_filename)
                                movie_filename = ''
                            elif media_extension in ["mp4","m4v","mov","avi","mpg","mpeg","wmv"]:
                                # Create Poster Image
                                pi_dir = imageDir + '/poster_images'
                                pi_path = pi_dir + '/' + os.path.splitext(os.path.basename(image_filename))[0] + ".jpeg"
                                if not os.path.exists(pi_dir):
                                    os.mkdir(pi_dir)
                                    print(bcolors.HEADER + 'Poster Image Directory Created')
                                if not os.path.exists(pi_path):
                                    create_poster_image(image_filename, pi_path)
                                    if not os.path.exists(pi_path):
                                        print(bcolors.FAIL + 'Failed to make poster image for video' + bcolors.ENDC)
                                        sys.exit(0)
                                
                                # Set Filenames
                                movie_filename = image_filename
                                image_filename = pi_path

                                # Temp insert picture for size calculation
                                picture_ph = picture.insert_picture(image_filename)
                            else:
                                raise Exception('Unsupported Image')
                                

                            available_width = picture_ph.width
                            available_height = picture_ph.height

                            image_width, image_height = picture_ph.image.size

                            placeholder_aspect_ratio = float(available_width) / float(available_height)
                            image_aspect_ratio = float(image_width) / float(image_height)

                            picture_ph.crop_top = 0
                            picture_ph.crop_left = 0
                            picture_ph.crop_bottom = 0
                            picture_ph.crop_right = 0

                            # ---if the placeholder is "wider" in aspect, shrink the picture width while
                            # ---maintaining the image aspect ratio
                            if placeholder_aspect_ratio > image_aspect_ratio:
                                picture_ph.width = int(image_aspect_ratio * available_height)
                                picture_ph.height = int(available_height)
                                picture_ph.top = picture_top
                                picture_ph.left = picture_left + int((available_width - picture_ph.width)/2)
                            # ---otherwise shrink the height
                            else:
                                picture_ph.height = int(available_width/image_aspect_ratio)
                                picture_ph.width = int(available_width)
                                picture_ph.top = picture_top + int((available_height - picture_ph.height)/2)
                                picture_ph.left = picture_left

                            media_px_width =  round(picture_ph.width / emus_per_inch * px_per_inch)
                            media_px_height =  round(picture_ph.height / emus_per_inch * px_per_inch) 
    
                            if movie_filename != '':
                                actual_media_width, actual_media_height = get_size(movie_filename)

                                if RESIZE_VIDEO and actual_media_width > media_px_width:

                                    video_dir = imageDir + '/scaled_videos'
                                    video_path = video_dir + '/' + os.path.splitext(os.path.basename(image_filename))[0] + ".mp4"
                                    
                                    if not os.path.exists(video_dir):
                                        os.mkdir(video_dir)
                                        print(bcolors.HEADER + 'Scaled Video Directory Created')
                                    if not os.path.exists(video_path):
                                        print("Resizing video...")    
                                        # Convert to divisible by 2 numbers for h264
                                        media_px_width = math.ceil(media_px_width/2)*2
                                        media_px_height = math.ceil(media_px_height/2)*2

                                        create_scaled_video(movie_filename, video_path, media_px_width, media_px_height)

                                        if not os.path.exists(video_path):
                                            print(bcolors.FAIL + 'Failed to make scaled video' + bcolors.ENDC)
                                            sys.exit(0)

                                    movie_filename = video_path
                                    media_extension = os.path.splitext(os.path.basename(movie_filename))[1].replace('.','')

                                # Set Mime Type
                                if media_extension in ["mp4","m4v"]:
                                    mtype = 'video/mp4'
                                elif media_extension in ["mov"]:
                                    mtype = 'video/quicktime'
                                elif media_extension in ["mpg","mpeg"]:
                                    mtype = 'video/mpeg'
                                elif media_extension in ["wmv"]:
                                    mtype = 'video/x-ms-wmv'
                                else:
                                    print(bcolors.FAIL + 'Media not supported: ' + str(media_extension) + bcolors.ENDC)
                                    sys.exit(0)

                                # Add movie
                                movie_shape = slide.shapes.add_movie(movie_filename, 0, 0, Inches(3), Inches(5), mime_type = mtype, poster_frame_image = image_filename)
                                
                                if AUTO_PLAY_VIDEO:
                                    # Autoplay video hack
                                    tree = movie_shape._element.getparent().getparent().getnext().getnext()
                                    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
                                    timing.set('delay', '0')

                                    # Mute
                                    medianode = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cMediaNode'][0]
                                    medianode.set('vol','0')

                                    # Loop
                                    repeat = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cTn'][1]
                                    repeat.set('repeatCount', 'indefinite')

                                # Match placeholder calculated size    
                                movie_shape.height = picture_ph.height
                                movie_shape.width = picture_ph.width
                                movie_shape.top = picture_ph.top
                                movie_shape.left = picture_ph.left

                                # Apply Border
                                # outline = movie_shape.line
                                # outline.color.rgb = RGBColor(255, 255, 255)
                                # outline.color.brightness = 1
                                # outline.width = Pt(4.5)

                                # Doesn't work
                                # picture_ph_outline = picture_ph.line
                                # outline.color.rgb = picture_ph_outline.color.rgb
                                # outline.color.brightness = picture_ph_outline.brightness
                                # outline.width = picture_ph_outline.width

                                # Remove temporary image in placeholder
                                if REMOVE_BLANK_IMAGE_PLACEHOLDER:
                                    print(bcolors.WARNING + 'Removing image placeholder from slide...' + bcolors.ENDC)
                                    picture_ph._element.getparent().remove(picture_ph._element)

                        else:
                            picture_ph = picture.insert_picture(image_filename)
                    except Exception as e:
                        print(bcolors.WARNING + 'Add Image Placeholder Error...' + bcolors.ENDC)
                        log_error(e)

                        if REMOVE_BLANK_IMAGE_PLACEHOLDER:
                            print(bcolors.WARNING + 'Removing image placeholder from slide...' + bcolors.ENDC)
                            picture._element.getparent().remove(picture._element)
                        elif MISSING_IMAGE_REPLACEMENT != "":
                            picture.insert_picture(MISSING_IMAGE_REPLACEMENT)
                        else:
                            print(bcolors.WARNING + 'Missing image and no replacement provided!' + bcolors.ENDC)
                        
                        warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " Image: " + image_filename + " - Bad Image")
                
                if QUOTE_COL != '':
                    if row[QUOTE_COL] != "":
                        quote.text = '"' + row[QUOTE_COL] + '"'
                        if MEASURE_QUOTE_HEIGHT:
                            quote_lines = int(quote_to_height.get(row[QUOTE_COL], 0))
                            if quote_lines == 0:
                                print(bcolors.WARNING + 'Quote Height Missing' + bcolors.ENDC)
                                print(row[QUOTE_COL])
                                print(bcolors.FAIL + 'Recalculate Quote Heights...' + bcolors.ENDC)
                                sys.exit(0)
                            print(bcolors.HEADER + 'Quote Lines needed: ' + str(quote_lines) + bcolors.ENDC)
                            quote_top = quote.top
                            quote_left = quote.left
                            quote_width = quote.width
                            quote_height = quote.height
                            if quote_lines <= 2:
                                print(bcolors.WARNING + 'Moving Content Down!' + bcolors.ENDC)
                                sub_top = subtext.top
                                sub_left = subtext.left
                                sub_width = subtext.width
                                sub_height = subtext.height

                                name_top = name.top
                                name_left = name.left
                                name_width = name.width
                                name_height = name.height

                                picture_top = picture_ph.top
                                picture_left = picture_ph.left
                                picture_width = picture_ph.width
                                picture_height = picture_ph.height

                                sub_offset = ((3 - quote_lines) * Pt(SUBTEXT_FONT_SIZE - 10))

                                subtext.top = sub_top + sub_offset
                                subtext.left = sub_left
                                subtext.width = sub_width
                                subtext.height = sub_height

                                name.top = name_top + sub_offset
                                name.left = name_left
                                name.width = name_width
                                name.height = name_height

                                picture_ph.top = picture_top + sub_offset
                                picture_ph.left = picture_left
                                picture_ph.width = picture_width
                                picture_ph.height = picture_height

                                quote.top = quote_top + sub_offset
                                quote.left = quote_left
                                quote.width = quote_width
                                quote.height = quote_height

                            if (quote_top + quote_lines * Inches(0.5) > Inches(prs_height) - Inches(0.5)): #Pt(2*SUBTEXT_FONT_SIZE)):
                                print(bcolors.FAIL + 'Too close to edge!' + bcolors.ENDC)
                                # sys.exit(0)
                                warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Too Close to Edge")
                            else:
                                print((Inches(prs_height) - (quote_top + quote_lines * Inches(0.5)))/emus_per_inch)
            except Exception as e:
                print(bcolors.WARNING + 'Student Slide Placeholder Error...' + bcolors.ENDC)
                print(bcolors.FAIL + str(e) + bcolors.ENDC)
                print("Current Layout:" + str(cur_layout))
                log_error(e)
                placeholder_error(slide)
            
            # ----------------------------------------------------------------------------------------
            # Create play button with attached audio
            # call_time('Audio Button')
            if not is_missing_audio:
                add_audio_button(slide, audio_filename, name)
            
            # call_time('Slide Timing')
            if AUTO_ADVANCE_SLIDE:
                if not is_missing_audio:
                    
                    # Auto advance slide hack
                    audio_dur = get_length(audio_filename)
                    # print('Audio duration: ' + str(audio_dur))
                    slide_dur = int((audio_dur + SLIDE_GAP)*1000)

                    if slide_dur < DEFAULT_SLIDE_DURATION:
                        # warnings_list.append(str(slide_count) + ": " + row[FULLNAME_COL] + " - Slide duration warning, setting default " + "[" + str(round(audio_dur,2)) + "|" + str(round(slide_dur/1000,2)) + "|" + str(round(DEFAULT_SLIDE_DURATION/1000,2))+ "]")
                        slide_dur = DEFAULT_SLIDE_DURATION
                    # print('Slide duration: ' + str(slide_dur))

                    print('Adding slide duration: ' + str(slide_dur/1000))

                    add_slide_duration(slide, slide_dur)
                else:
                    print('Adding default slide duration: ' + str(DEFAULT_SLIDE_DURATION/1000))
                    add_slide_duration(slide, DEFAULT_SLIDE_DURATION)
            
            # Prepare for next slide
            cur_layout += 1 #not cur_layout
            if cur_layout > MULTI_LAYOUT_END_NUM: # max
                cur_layout = STUDENT_SLIDE_LAYOUT_NUM

            line_count += 1
            slide_count += 1
            is_missing_audio = False
            is_missing_image = False
            movie_filename = ''
            image_filename = ''

            # Prevent server spamming if we just downloaded some files
            # call_time('Download Delay')
            if files_downloaded:
                # delay = round(random.uniform(0.2,1.0), 2) # random delay
                # print(bcolors.WARNING + 'Delaying next download (' + str(delay) + 'sec)...' + bcolors.ENDC)
                # sleep(delay) # sleep...
                files_downloaded = False

    # End CSV loop
    # call_time('Last Slide')
    csv_file.close()
    print(bcolors.OKBLUE + '--------------------' + bcolors.ENDC)
    print(bcolors.HEADER + '\nProccessed ' + str(line_count) + ' lines.')
    print(bcolors.HEADER + 'Created ' + str(slide_count) + ' slides.')
    print(bcolors.HEADER + 'Created ' + str(title_slide_count) + ' title slides at:' + str(list(title_slides_list)) + '\n')

    if AUTO_ADVANCE_SLIDE:
        print (bcolors.WARNING + "Total slide time: " + millis_to_hms(total_slide_time))
 
    print(bcolors.WARNING + 'Total missing audio: ' + str(total_missing_audio))
    print('Total missing images: ' + str(total_missing_photos) + '\n' + bcolors.ENDC)

    print(bcolors.WARNING + 'Slide Warnings:')
    for warning in warnings_list:
        print(warning)
    print('\n' + bcolors.ENDC)

    # Save Final File
    # call_time('Save PPT')
    if SPLIT_OUTPUT:
        if SPLIT_BY_NUM:
            save_presentation(prs, "[" + str(split_total*SLIDES_PER_PPT+1) + "-" + str(slide_count-1) + "]")
        else:
            # save_presentation(prs, "[" + str(cur_master) + "]")
            save_presentation(prs, "[" + str(cur_master) + "] " + title_text)
    else:
        save_presentation(prs, "")

    print(bcolors.OKGREEN + '------- Done -------' + bcolors.ENDC)
    end_time = time.perf_counter() - start_time
    print(bcolors.OKBLUE + 'Total Run Time: ' + millis_to_hms(end_time * 1000))